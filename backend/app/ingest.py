import re
import os
from typing import List
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl

from app.vector_store import add_embeddings
from app.config import EmbeddingClient

CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
UPLOAD_DIR = "data/raw"


def read_pdf_text(file_path: str) -> str:
    reader = PdfReader(file_path)
    text_parts = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        text_parts.append(page_text)
    return "\n".join(text_parts)


def read_text_file(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()


def read_docx_text(file_path: str) -> str:
    doc = Document(file_path)
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def read_pptx_text(file_path: str) -> str:
    prs = Presentation(file_path)
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    slide_lines.append(text)
        if slide_lines:
            parts.append(f"Slide {slide_idx}:")
            parts.append("\n".join(slide_lines))
            parts.append("")

    return "\n".join(parts)


def read_xlsx_text(file_path: str) -> str:
    sheets = pd.read_excel(file_path, sheet_name=None)
    parts: list[str] = []

    week_pattern = re.compile(r"Week\s*\d+", re.IGNORECASE)
    date_pattern = re.compile(r"\d{1,2}\s\w+\s\d{4}")

    for sheet_name, df in sheets.items():

        df = df.dropna(how="all", axis=0).dropna(how="all", axis=1)

        if df.empty:
            continue

        df = df.astype(str)

        parts.append(f"Sheet: {sheet_name}")

        for _, row in df.iterrows():
            vals = [v.strip() for v in row.tolist()]

            joined = " ".join(vals).lower()
            if joined.count("nan") > 4:
                continue

            if vals.count("") > len(vals) * 0.7:
                continue

            row_str = ", ".join(vals)

            if (
                "week" in joined
                or date_pattern.search(row_str)
            ):
                parts.append(row_str)

        parts.append("")

    return "\n".join(parts)


def chunk_text(
    text: str,
    chunk_size: int = CHUNK_SIZE,
    overlap: int = CHUNK_OVERLAP,
) -> List[str]:
    chunk_size = int(chunk_size)
    overlap = int(overlap)

    if chunk_size <= 0:
        raise ValueError("chunk_size must be > 0")
    if overlap < 0:
        raise ValueError("overlap must be >= 0")
    if overlap >= chunk_size:
        overlap = max(0, chunk_size // 5)

    chunks = []
    start = 0
    length = len(text)

    while start < length:
        end = min(start + chunk_size, length)
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        if end == length:
            break
        start = end - overlap
    return chunks


def ingest_file(
    file_path: str,
    doc_name: str,
    chunk_size: int = CHUNK_SIZE,
    chunk_overlap: int = CHUNK_OVERLAP,
) -> int:
    if not os.path.exists(file_path):
        raise FileNotFoundError(file_path)

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        full_text = read_pdf_text(file_path)
    elif ext in {".txt", ".csv"}:
        full_text = read_text_file(file_path)
    elif ext == ".docx":
        full_text = read_docx_text(file_path)
    elif ext == ".pptx":
        full_text = read_pptx_text(file_path)
    elif ext == ".xlsx":
        full_text = read_xlsx_text(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    chunks = chunk_text(full_text, chunk_size=chunk_size,
                        overlap=chunk_overlap)
    if not chunks:
        return 0

    embed_client = EmbeddingClient()
    embeddings = embed_client.embed_documents(chunks)

    add_embeddings(chunks, embeddings, doc_name=doc_name)
    return len(chunks)
